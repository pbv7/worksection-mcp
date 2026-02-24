"""File-related MCP tools."""

import base64
import io
import mimetypes

from worksection_mcp.cache import FileCache
from worksection_mcp.client import WorksectionClient
from worksection_mcp.mcp_protocols import ToolRegistrar


def register_file_tools(
    mcp: ToolRegistrar, client: WorksectionClient, file_cache: FileCache | None = None
) -> None:
    """Register file-related tools with the MCP server."""

    async def _get_all_task_attachments_internal(task_id: str) -> dict:
        """Internal helper to get all attachments from a task and its comments.

        This is used by both get_all_task_attachments and list_image_attachments
        to avoid the issue where MCP-decorated functions cannot call each other.
        """
        task_data = await client.get_task(task_id=task_id, extra="files")
        comments_data = await client.get_comments(task_id=task_id, extra="files")

        task_files = []
        if isinstance(task_data, dict):
            # Files can be in task_data["data"]["files"] or task_data["files"]
            data = task_data.get("data", task_data)
            task_files = data.get("files", [])

        comment_files = [
            {
                **file,
                "comment_id": comment.get("id"),
                "comment_text": comment.get("text", "")[:100],  # Preview
                "comment_author": comment.get("user_from", {}).get("name"),
            }
            for comment in (
                comments_data.get("data", [])
                if isinstance(comments_data, dict) and "data" in comments_data
                else []
            )
            for file in comment.get("files", [])
        ]

        return {
            "task_id": task_id,
            "task_files": task_files,
            "comment_files": comment_files,
            "total_files": len(task_files) + len(comment_files),
        }

    @mcp.tool()
    async def get_project_files(
        project_id: str | None = None,
        task_id: str | None = None,
    ) -> dict:
        """List all files in a project or task.

        Includes Files section uploads, description attachments, and comment files.
        Unlike get_task_files (which uses extra=files), this calls the dedicated
        get_files API endpoint for broader coverage.

        Args:
            project_id: Project ID (at least one of project_id or task_id required)
            task_id: Task ID (at least one of project_id or task_id required)

        Returns:
            List of files with metadata:
            - id: File ID (use with download_file)
            - name: File name
            - size: File size
            - date_added: Upload date
        """
        if not project_id and not task_id:
            return {"error": "At least one of project_id or task_id is required."}
        return await client.get_files(project_id=project_id, task_id=task_id)

    @mcp.tool()
    async def get_task_files(task_id: str) -> dict:
        """Get all files attached to a task.

        Args:
            task_id: The task ID

        Returns:
            List of attached files with metadata:
            - id: File ID (use with download_file)
            - name: File name
            - size: File size in bytes
            - type: MIME type
            - date_added: Upload date
            - user: Who uploaded the file
        """
        return await client.get_task(task_id=task_id, extra="files")

    @mcp.tool()
    async def get_all_task_attachments(task_id: str) -> dict:
        """Get all attachments from a task and its comments.

        This collects files from both the task itself and all comments,
        useful for getting a complete picture of all visual content.

        Args:
            task_id: The task ID

        Returns:
            All attachments organized by source:
            - task_files: Files attached directly to task
            - comment_files: Files from comments (with comment context)
            - total_files: Total attachment count
        """
        return await _get_all_task_attachments_internal(task_id)

    @mcp.tool()
    async def download_file(file_id: str) -> dict:
        """Download a file and return its metadata.

        The file is cached locally for subsequent access via MCP resources.

        Args:
            file_id: The file ID to download

        Returns:
            File metadata and cache information:
            - file_id: The file ID
            - cached: Whether file was cached successfully
            - resource_uri: MCP resource URI to access the file
            - size_bytes: File size
            - mime_type: Detected MIME type
        """
        content = await client.download_file(file_id=file_id)

        # Try to detect mime type (will be refined if file_cache available)
        mime_type = "application/octet-stream"

        # If file cache is available, save the file
        cache_path = None
        if file_cache:
            cache_path = await file_cache.save(file_id, content)
            mime_type = mimetypes.guess_type(str(cache_path))[0] or mime_type

        return {
            "file_id": file_id,
            "cached": cache_path is not None,
            "resource_uri": f"worksection://file/{file_id}",
            "size_bytes": len(content),
            "mime_type": mime_type,
        }

    def _detect_mime_type(content: bytes) -> str:
        """Detect MIME type from file content magic bytes."""
        # Common file signatures
        signatures = {
            b"\x89PNG\r\n\x1a\n": "image/png",
            b"\xff\xd8\xff": "image/jpeg",
            b"GIF87a": "image/gif",
            b"GIF89a": "image/gif",
            b"RIFF": "image/webp",  # WebP starts with RIFF....WEBP
            b"BM": "image/bmp",
            b"PK\x03\x04": "application/vnd.openxmlformats-officedocument",  # DOCX, XLSX, PPTX
            b"%PDF": "application/pdf",
        }

        for sig, mime in signatures.items():
            if content.startswith(sig):
                # Special handling for Office formats
                if mime == "application/vnd.openxmlformats-officedocument":
                    # Check for specific Office format markers in the ZIP
                    if b"word/" in content[:2000]:
                        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    if b"xl/" in content[:2000]:
                        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    if b"ppt/" in content[:2000]:
                        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    return mime
                # WebP needs additional check
                if sig == b"RIFF" and b"WEBP" in content[:12]:
                    return "image/webp"
                if sig == b"RIFF":
                    continue  # Not WebP, skip
                return mime

        # Heuristic text detection for plain UTF-8 content without magic bytes.
        # This keeps binary payloads in octet-stream while allowing readable text files.
        try:
            decoded = content.decode("utf-8")
            if decoded and all((ch.isprintable() or ch in "\r\n\t") for ch in decoded):
                return "text/plain"
        except UnicodeDecodeError:
            pass

        return "application/octet-stream"

    @mcp.tool()
    async def get_file_as_base64(file_id: str) -> dict:
        """Download a file and return it as base64 encoded string.

        Useful for directly including small files (like images) in responses.
        For images, returns data URL format ready for display.
        For larger files, use download_file and access via MCP resources.

        Args:
            file_id: The file ID to download

        Returns:
            File data as base64:
            - file_id: The file ID
            - base64_content: Base64 encoded file content
            - size_bytes: Original file size
            - mime_type: Detected MIME type
            - data_url: For images, a data URL ready for display
            - is_image: Whether this is an image file
        """
        content = await client.download_file(file_id=file_id)

        # Detect MIME type from content
        mime_type = _detect_mime_type(content)

        # Encode as base64
        base64_content = base64.b64encode(content).decode("utf-8")

        # Check if it's an image
        is_image = mime_type.startswith("image/")

        result = {
            "file_id": file_id,
            "base64_content": base64_content,
            "size_bytes": len(content),
            "mime_type": mime_type,
            "is_image": is_image,
        }

        # For images, provide data URL for easy display
        if is_image:
            result["data_url"] = f"data:{mime_type};base64,{base64_content}"

        return result

    @mcp.tool()
    async def list_image_attachments(
        task_id: str | None = None,
        project_id: str | None = None,
    ) -> dict:
        """List all image attachments for a task or project.

        Filters to only image files (jpg, png, gif, webp, etc.)
        for easy access to screenshots and visual content.

        Exactly one of task_id or project_id must be provided.

        Args:
            task_id: The task ID (gets images from task and its comments)
            project_id: The project ID (gets images from project files)

        Returns:
            List of image files:
            - images: Array of image file metadata
            - Each with: id, name, resource_uri, source (task/comment/project)
        """
        if task_id and project_id:
            raise ValueError("Provide either task_id or project_id, not both")
        if not task_id and not project_id:
            raise ValueError("Either task_id or project_id is required")

        image_extensions = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".svg"}
        images = []

        if project_id:
            # Get project files and filter to images
            files_data = await client.get_files(project_id=project_id)
            if isinstance(files_data, dict) and "data" in files_data:
                for file in files_data["data"]:
                    name = file.get("name", "")
                    if any(name.lower().endswith(ext) for ext in image_extensions):
                        images.append(
                            {
                                "id": file.get("id"),
                                "name": name,
                                "resource_uri": f"worksection://file/{file.get('id')}",
                                "source": "project",
                            }
                        )
            return {
                "project_id": project_id,
                "images": images,
                "image_count": len(images),
            }

        # task_id is guaranteed non-None here (validated above, project_id returned early)
        all_attachments = await _get_all_task_attachments_internal(str(task_id))

        # Check task files
        for file in all_attachments.get("task_files", []):
            name = file.get("name", "")
            if any(name.lower().endswith(ext) for ext in image_extensions):
                images.append(
                    {
                        "id": file.get("id"),
                        "name": name,
                        "resource_uri": f"worksection://file/{file.get('id')}",
                        "source": "task",
                    }
                )

        # Check comment files
        for file in all_attachments.get("comment_files", []):
            name = file.get("name", "")
            if any(name.lower().endswith(ext) for ext in image_extensions):
                images.append(
                    {
                        "id": file.get("id"),
                        "name": name,
                        "resource_uri": f"worksection://file/{file.get('id')}",
                        "source": "comment",
                        "comment_preview": file.get("comment_text"),
                    }
                )

        return {
            "task_id": task_id,
            "images": images,
            "image_count": len(images),
        }

    def _extract_text_from_docx(content: bytes) -> str:
        """Extract text from DOCX file."""
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            return f"[Error extracting DOCX text: {e}]"

    def _extract_text_from_xlsx(content: bytes) -> str:
        """Extract text from XLSX file."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            result = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                result.append(f"=== Sheet: {sheet_name} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):
                        result.append("\t".join(row_values))
            wb.close()
            return "\n".join(result)
        except Exception as e:
            return f"[Error extracting XLSX text: {e}]"

    def _extract_text_from_pptx(content: bytes) -> str:
        """Extract text from PPTX file."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            result = []
            for slide_num, slide in enumerate(prs.slides, 1):
                result.append(f"=== Slide {slide_num} ===")
                for shape in slide.shapes:
                    text = getattr(shape, "text", None)
                    if isinstance(text, str) and text.strip():
                        result.append(text)
            return "\n\n".join(result)
        except Exception as e:
            return f"[Error extracting PPTX text: {e}]"

    def _extract_text_from_pdf(content: bytes) -> str:
        """Extract text from PDF file."""
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(content))
            result = []
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    result.append(f"=== Page {page_num} ===")
                    result.append(text)
            return "\n\n".join(result)
        except Exception as e:
            return f"[Error extracting PDF text: {e}]"

    def _extract_text_content(content: bytes, mime_type: str) -> tuple[str, str]:
        """Extract text from file based on MIME type.

        Returns:
            Tuple of (extracted_text, content_type)
            content_type is one of: text, document, spreadsheet, presentation, image, binary
        """
        # Text files
        if mime_type.startswith("text/") or mime_type in [
            "application/json",
            "application/xml",
            "application/javascript",
        ]:
            try:
                return content.decode("utf-8"), "text"
            except UnicodeDecodeError:
                try:
                    return content.decode("latin-1"), "text"
                except Exception:
                    return "[Unable to decode text content]", "binary"

        # Word documents
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _extract_text_from_docx(content), "document"

        # Excel spreadsheets
        if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return _extract_text_from_xlsx(content), "spreadsheet"

        # PowerPoint presentations
        if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return _extract_text_from_pptx(content), "presentation"

        # PDF
        if mime_type == "application/pdf":
            return _extract_text_from_pdf(content), "document"

        # Images - can't extract text, but note they can be viewed
        if mime_type.startswith("image/"):
            return "[Image file - use get_file_as_base64 to view]", "image"

        return "[Binary file - cannot extract text]", "binary"

    @mcp.tool()
    async def get_file_content(file_id: str) -> dict:
        """Download a file and extract its readable content.

        Supports text extraction from:
        - Text files (.txt, .md, .json, .xml, etc.)
        - Word documents (.docx)
        - Excel spreadsheets (.xlsx)
        - PowerPoint presentations (.pptx)
        - PDF documents (.pdf)

        For images, use get_file_as_base64 instead.

        Args:
            file_id: The file ID to download

        Returns:
            Extracted file content:
            - file_id: The file ID
            - mime_type: Detected MIME type
            - content_type: Type category (text, document, spreadsheet, presentation, image, binary)
            - text_content: Extracted text (if available)
            - size_bytes: Original file size
            - is_readable: Whether text was successfully extracted
        """
        content = await client.download_file(file_id=file_id)

        # Detect MIME type
        mime_type = _detect_mime_type(content)

        # Extract text content
        text_content, content_type = _extract_text_content(content, mime_type)

        is_readable = content_type not in ("image", "binary")

        return {
            "file_id": file_id,
            "mime_type": mime_type,
            "content_type": content_type,
            "text_content": text_content,
            "size_bytes": len(content),
            "is_readable": is_readable,
        }
